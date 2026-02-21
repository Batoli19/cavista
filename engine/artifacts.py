from __future__ import annotations

import uuid
from pathlib import Path
from typing import Any, Dict, List


GENERATED_DIR = Path(__file__).resolve().parent.parent / "generated"


class MissingExportDependencyError(Exception):
    pass


def _ensure_dir() -> Path:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    return GENERATED_DIR


def _safe_name(name: str, default: str) -> str:
    cleaned = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in (name or "").strip())
    cleaned = cleaned.strip("_")
    return cleaned or default


def _file_meta(path: Path, file_type: str) -> Dict[str, Any]:
    fid = path.stem
    return {
        "id": fid,
        "type": file_type,
        "name": path.name,
        "url": f"/download/{fid}",
        "size": path.stat().st_size if path.exists() else 0,
        "path": str(path),
    }


def _research_sections(research_obj: Dict[str, Any]) -> Dict[str, Any]:
    topic = research_obj.get("topic", "Research")
    summary = research_obj.get("summary", "")
    key_points = research_obj.get("key_points", []) or []
    data_points = research_obj.get("data_points", []) or []
    sources = research_obj.get("sources", []) or []
    return {
        "topic": topic,
        "summary": summary,
        "key_points": key_points[:10],
        "data_points": data_points[:100],
        "sources": sources[:25],
    }


def export_docx_from_research(research_obj: Dict[str, Any]) -> Dict[str, Any]:
    try:
        from docx import Document
    except ModuleNotFoundError as e:
        raise MissingExportDependencyError("python-docx") from e

    s = _research_sections(research_obj)
    out_dir = _ensure_dir()
    fid = str(uuid.uuid4())[:12]
    name = _safe_name(s["topic"], "research")
    out_path = out_dir / f"{fid}_{name}.docx"

    doc = Document()
    doc.add_heading(s["topic"], level=1)
    doc.add_heading("Executive Summary", level=2)
    doc.add_paragraph(s["summary"] or "No summary available.")

    doc.add_heading("Key Points", level=2)
    if s["key_points"]:
        for p in s["key_points"]:
            doc.add_paragraph(str(p), style="List Bullet")
    else:
        doc.add_paragraph("No key points available.")

    doc.add_heading("Data Points", level=2)
    if s["data_points"]:
        for d in s["data_points"]:
            if isinstance(d, dict):
                line = ", ".join([f"{k}: {v}" for k, v in d.items()])
            else:
                line = str(d)
            doc.add_paragraph(line, style="List Bullet")
    else:
        doc.add_paragraph("No structured data points available.")

    doc.add_heading("Sources", level=2)
    if s["sources"]:
        for src in s["sources"]:
            title = src.get("title", "Source")
            domain = src.get("domain", "")
            note = src.get("note", "")
            url = src.get("url", "")
            doc.add_paragraph(f"{title} ({domain}) - {note}")
            if url:
                doc.add_paragraph(url)
    else:
        doc.add_paragraph("No sources available.")

    doc.save(str(out_path))
    return _file_meta(out_path, "docx")


def export_pptx_from_research(research_obj: Dict[str, Any]) -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as e:
        raise MissingExportDependencyError("python-pptx") from e

    s = _research_sections(research_obj)
    out_dir = _ensure_dir()
    fid = str(uuid.uuid4())[:12]
    name = _safe_name(s["topic"], "research")
    out_path = out_dir / f"{fid}_{name}.pptx"

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = s["topic"]
    slide.placeholders[1].text = "Research Brief"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.shapes.placeholders[1].text = s["summary"] or "No summary available."

    points = s["key_points"] or ["No key points available."]
    for i, point in enumerate(points[:5], start=1):
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = f"Insight {i}"
        sl.shapes.placeholders[1].text = str(point)

    src_slide = prs.slides.add_slide(prs.slide_layouts[1])
    src_slide.shapes.title.text = "Sources"
    tf = src_slide.shapes.placeholders[1].text_frame
    tf.clear()
    if s["sources"]:
        for src in s["sources"][:8]:
            p = tf.add_paragraph()
            p.text = f"{src.get('title', 'Source')} ({src.get('domain', '')})"
    else:
        tf.text = "No sources available."

    prs.save(str(out_path))
    return _file_meta(out_path, "pptx")


def export_xlsx_from_research(research_obj: Dict[str, Any]) -> Dict[str, Any]:
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference
    except ModuleNotFoundError as e:
        raise MissingExportDependencyError("openpyxl") from e

    s = _research_sections(research_obj)
    out_dir = _ensure_dir()
    fid = str(uuid.uuid4())[:12]
    name = _safe_name(s["topic"], "research")
    out_path = out_dir / f"{fid}_{name}.xlsx"

    wb = Workbook()
    ws_summary = wb.active
    ws_summary.title = "Summary"
    ws_summary.append(["Topic", s["topic"]])
    ws_summary.append(["Executive Summary", s["summary"] or "No summary available."])
    ws_summary.append([])
    ws_summary.append(["Key Points"])
    for p in s["key_points"] or ["No key points available."]:
        ws_summary.append([str(p)])

    ws_data = wb.create_sheet("Data")
    ws_data.append(["label", "value"])
    numeric_rows = 0
    for d in s["data_points"]:
        if isinstance(d, dict) and "label" in d and "value" in d:
            ws_data.append([str(d["label"]), float(d["value"]) if str(d["value"]).replace(".", "", 1).isdigit() else 0])
            numeric_rows += 1
        elif isinstance(d, dict):
            first_k = next(iter(d.keys()), "item")
            first_v = d.get(first_k, "")
            val = float(first_v) if str(first_v).replace(".", "", 1).isdigit() else 0
            ws_data.append([first_k, val])
            numeric_rows += 1

    if numeric_rows >= 2:
        chart = BarChart()
        chart.title = "Data Overview"
        data = Reference(ws_data, min_col=2, min_row=1, max_row=numeric_rows + 1)
        cats = Reference(ws_data, min_col=1, min_row=2, max_row=numeric_rows + 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        ws_data.add_chart(chart, "D2")

    ws_sources = wb.create_sheet("Sources")
    ws_sources.append(["title", "domain", "note", "url"])
    for src in s["sources"]:
        ws_sources.append([src.get("title", ""), src.get("domain", ""), src.get("note", ""), src.get("url", "")])

    wb.save(str(out_path))
    return _file_meta(out_path, "xlsx")
